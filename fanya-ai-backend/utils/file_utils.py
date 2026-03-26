import os
import re
import threading
from typing import Optional

_converter = None
_converter_lock = threading.Lock()
_converter_ready = threading.Event()


def allowed_file(filename):
    ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx'}
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def _init_converter():
    global _converter
    if _converter is not None:
        _converter_ready.set()
        return
    with _converter_lock:
        if _converter is not None:
            _converter_ready.set()
            return
        try:
            from docling.document_converter import DocumentConverter
            _converter = DocumentConverter()
            print("[Docling] DocumentConverter 初始化成功")
        except Exception as e:
            print(f"[Docling] 初始化失败，将降级为 PyPDF2: {e}")
            _converter = None
        finally:
            _converter_ready.set()


def start_converter_warmup():
    if not _converter_ready.is_set():
        threading.Thread(target=_init_converter, daemon=True, name="docling-warmup").start()


def _extract_pdf_pypdf2(file_path):
    import PyPDF2
    pages = []
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            page_text = (page.extract_text() or "").strip()
            # 保留空页占位，确保 content_list[i] 始终对应 PDF 第 i+1 页
            pages.append(page_text)
    return pages


def _is_text_quality_acceptable(pages):
    """
    判断 PyPDF2 提取文本是否可用于 RAG。
    某些扫描件或复杂字体 PDF 会提取到大量乱码，字符数虽然够，但检索不可用。
    """
    if not pages:
        return False

    text = "\n".join([p for p in pages if p]).strip()
    if not text:
        return False

    total_len = len(text)
    if total_len < 80:
        return False

    readable_chars = re.findall(r"[\u4e00-\u9fffA-Za-z0-9]", text)
    readable_ratio = len(readable_chars) / max(total_len, 1)

    # replacement char 出现通常意味着解码/字体映射失败
    replacement_count = text.count("\ufffd")
    replacement_ratio = replacement_count / max(total_len, 1)

    # 最少需要一定“词样式”片段，避免全是碎符号
    token_like = re.findall(r"[\u4e00-\u9fff]{2,}|[A-Za-z]{3,}", text)

    return readable_ratio >= 0.35 and replacement_ratio < 0.02 and len(token_like) >= 3


def _extract_pdf_pypdfium2(file_path):
    """使用 pypdfium2 按页提取文本，保留页码对应关系（比 PyPDF2 更准确）。"""
    try:
        import pypdfium2 as pdfium
        doc = pdfium.PdfDocument(file_path)
        pages = []
        for i in range(len(doc)):
            page = doc[i]
            textpage = page.get_textpage()
            text = (textpage.get_text_range() or "").strip()
            pages.append(text)
        return pages
    except Exception as e:
        print(f"[pypdfium2] 文本提取失败: {e}")
        return []


def _extract_pdf_docling(file_path):
    """使用 Docling 尝试按页提取文本；若无法按页则返回空列表（避免段落错位）。"""
    if not _converter_ready.is_set() or _converter is None:
        _init_converter()
    if _converter is None:
        return []
    try:
        result = _converter.convert(file_path)
        doc = result.document
        # 尝试按页获取文本（保留页码对应关系）
        if hasattr(doc, 'pages') and doc.pages:
            pages = []
            for page in doc.pages:
                page_text = ""
                if hasattr(page, 'export_to_markdown'):
                    page_text = (page.export_to_markdown() or "").strip()
                elif hasattr(page, 'text'):
                    page_text = (page.text or "").strip()
                pages.append(page_text)
            if any(pages):
                return pages
        # 无法按页获取时放弃，由下游 RapidOCR 按页处理
        return []
    except Exception as e:
        print(f"[Docling] 提取失败: {e}")
        return []


def _extract_pdf_rapidocr(file_path, max_pages=30):
    """
    轻量 OCR 兜底：当 Docling 不可用或失败时，使用 rapidocr 识别 PDF 页图像。
    保留空页占位（空字符串），确保返回列表与 PDF 页码一一对应。
    """
    try:
        import numpy as np
        import pypdfium2 as pdfium
        from rapidocr import RapidOCR
    except Exception as e:
        print(f"[RapidOCR] 依赖不可用: {e}")
        return []

    try:
        doc = pdfium.PdfDocument(file_path)
        ocr_engine = RapidOCR()
        pages = []
        page_count = min(len(doc), max_pages)

        for i in range(page_count):
            page = doc[i]
            pil_img = page.render(scale=2.0).to_pil()
            img = np.array(pil_img)
            ocr_res, _ = ocr_engine(img)
            if not ocr_res:
                pages.append("")  # 保留空页占位
                continue

            lines = []
            for item in ocr_res:
                if not isinstance(item, (list, tuple)) or len(item) < 2:
                    continue
                raw_text = item[1]
                if isinstance(raw_text, str):
                    text = raw_text.strip()
                elif isinstance(raw_text, (list, tuple)) and raw_text:
                    text = str(raw_text[0]).strip()
                else:
                    text = ""
                if text:
                    lines.append(text)

            pages.append(" ".join(lines) if lines else "")

        return pages
    except Exception as e:
        print(f"[RapidOCR] 识别失败: {e}")
        return []


def extract_text_from_pdf(file_path):
    """
    从PDF文件中按页提取文本。
    核心原则：返回列表的第 i 项必须对应 PDF 第 i+1 页，绝不跳过空页，
    以保证 aiScript[i].page == i+1 与 PDF 页码完全对应。
    """
    try:
        # 优先 PyPDF2（快速，已修正为保留空页）
        fast_result = _extract_pdf_pypdf2(file_path)
        total_chars = sum(len(p.strip()) for p in fast_result)
        if total_chars >= 100 and _is_text_quality_acceptable(fast_result):
            print(f"[PDF] PyPDF2 快速提取成功（{total_chars} 字符，质量通过）")
            return fast_result

        # 次选 pypdfium2（比 PyPDF2 更准确，同样按页返回）
        print(f"[PDF] PyPDF2 文本质量不足（{total_chars} 字符），尝试 pypdfium2")
        pypdfium2_result = _extract_pdf_pypdfium2(file_path)
        pypdfium2_chars = sum(len(p.strip()) for p in pypdfium2_result)
        if pypdfium2_chars >= 100 and _is_text_quality_acceptable(pypdfium2_result):
            print(f"[PDF] pypdfium2 提取成功（{pypdfium2_chars} 字符）")
            return pypdfium2_result

        # 再试 Docling（仅在能按页返回时使用，避免段落错位）
        print(f"[PDF] pypdfium2 文本质量不足（{pypdfium2_chars} 字符），尝试 Docling")
        docling_result = _extract_pdf_docling(file_path)
        if docling_result:
            print(f"[PDF] Docling 按页提取成功（{len(docling_result)} 页）")
            return docling_result

        # 最后兜底：RapidOCR 逐页 OCR（图片类 PDF），现已保留空页占位
        print("[PDF] Docling 未生效，尝试 RapidOCR 逐页 OCR 兜底")
        rapidocr_result = _extract_pdf_rapidocr(file_path)
        if rapidocr_result and _is_text_quality_acceptable(rapidocr_result):
            print("[PDF] RapidOCR 识别成功")
            return rapidocr_result

        print("[PDF] 所有提取方式均未命中有效文本，回退 PyPDF2 结果（含空页占位）")
        return fast_result
    except Exception as e:
        print(f"PDF解析失败: {e}")
        return []


def extract_text_from_ppt(file_path):
    """从PPT文件中提取文本"""
    try:
        from pptx import Presentation
        text = []
        prs = Presentation(file_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_text.append(shape.text)
            text.append(' '.join(slide_text))
        return text
    except Exception as e:
        print(f"PPT解析失败: {e}")
        return []


def _ppt_to_pdf_libreoffice(file_path: str, pdf_path: str) -> Optional[str]:
    """使用 LibreOffice 无头模式将 ppt/pptx 转为 pdf（不依赖 WPS）。"""
    import subprocess
    import shutil

    base_dir = os.path.dirname(os.path.abspath(file_path))
    candidates = []
    env_exe = os.environ.get("LIBREOFFICE_SOFFICE", "").strip()
    if env_exe:
        candidates.append(env_exe)
    if os.name == "nt":
        candidates.extend(
            [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            ]
        )
    candidates.append("soffice")

    soffice = None
    for c in candidates:
        if not c:
            continue
        if c == "soffice":
            soffice = shutil.which("soffice")
            if soffice:
                break
        elif os.path.isfile(c):
            soffice = c
            break

    if not soffice:
        print("[PPT2PDF] LibreOffice (soffice) 未找到，跳过。可安装 LibreOffice 或设置 LIBREOFFICE_SOFFICE")
        return None

    abs_in = os.path.abspath(file_path)
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", base_dir, abs_in],
            check=True,
            timeout=300,
            capture_output=True,
            text=True,
        )
    except Exception as e:
        print(f"[PPT2PDF] LibreOffice 转换失败: {e}")
        return None

    if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
        print(f"[PPT2PDF] LibreOffice 转换成功: {pdf_path}")
        return pdf_path
    print(f"[PPT2PDF] LibreOffice 未生成预期文件: {pdf_path}")
    return None


def _ppt_to_pdf_powerpoint_com(file_path: str, pdf_path: str) -> Optional[str]:
    """使用 Microsoft PowerPoint COM 导出 PDF（已安装 Office 时可用）。"""
    try:
        import comtypes.client
        import pythoncom
    except ImportError:
        return None

    abs_file_path = os.path.abspath(file_path)
    abs_pdf_path = os.path.abspath(pdf_path)
    pythoncom.CoInitialize()
    app = None
    pres = None
    try:
        app = comtypes.client.CreateObject("PowerPoint.Application")
        app.Visible = 1
        pres = app.Presentations.Open(abs_file_path, WithWindow=False)
        # 2 = ppFixedFormatTypePDF
        pres.ExportAsFixedFormat(abs_pdf_path, 2)
        pres.Close()
        app.Quit()
        if os.path.exists(abs_pdf_path) and os.path.getsize(abs_pdf_path) > 0:
            print("[PPT2PDF] Microsoft PowerPoint COM 转换成功")
            return pdf_path
        return None
    except Exception as e:
        print(f"[PPT2PDF] PowerPoint COM 失败: {e}")
        if pres:
            try:
                pres.Close()
            except Exception:
                pass
        if app:
            try:
                app.Quit()
            except Exception:
                pass
        return None
    finally:
        try:
            import pythoncom
            pythoncom.CoUninitialize()
        except Exception:
            pass


def ppt_to_pdf(file_path):
    """将 PPT/PPTX 转为 PDF：依次尝试 WPS 演示 COM → Microsoft PowerPoint COM（不使用 LibreOffice 等其它方式）。"""
    import time
    import subprocess

    print(f"开始PPT转PDF: {file_path}")

    base_dir = os.path.dirname(file_path)
    filename = os.path.basename(file_path)
    name_without_ext = os.path.splitext(filename)[0]
    pdf_path = os.path.join(base_dir, f"{name_without_ext}.pdf")

    print(f"目标PDF路径: {pdf_path}")

    def _kill_wps():
        for proc in ['wps.exe', 'wpp.exe', 'et.exe']:
            try:
                subprocess.run(['taskkill', '/F', '/IM', proc],
                               capture_output=True, timeout=5)
            except Exception:
                pass

    # ---------- 1) WPS ----------
    try:
        import comtypes.client
        import pythoncom

        pythoncom.CoInitialize()

        _kill_wps()
        time.sleep(0.5)
        print("[PPT2PDF] cleaned stale WPS processes")

        abs_file_path = os.path.abspath(file_path)
        abs_pdf_path = os.path.abspath(pdf_path)

        wps_com_names = [
            "KWPP.Application",
            "WPP.Application",
            "Kwpp.Application",
            "wpp.Application",
        ]

        for com_name in wps_com_names:
            app = None
            presentation = None
            try:
                print(f"[PPT2PDF] trying COM object: {com_name}")
                app = comtypes.client.CreateObject(com_name)
                app.Visible = 1
                presentation = app.Presentations.Open(abs_file_path)
                presentation.ExportAsFixedFormat(abs_pdf_path, 2)
                presentation.Close()
                app.Quit()
                if os.path.exists(abs_pdf_path) and os.path.getsize(abs_pdf_path) > 0:
                    print(f"[PPT2PDF] conversion success via {com_name}")
                    return pdf_path
            except Exception as wps_error:
                print(f"[PPT2PDF] {com_name} convert failed: {wps_error}")
                if presentation:
                    try:
                        presentation.Close()
                    except Exception:
                        pass
                if app:
                    try:
                        app.Quit()
                    except Exception:
                        pass
            finally:
                _kill_wps()
                time.sleep(0.3)
    except ImportError as import_error:
        print(f"[PPT2PDF] comtypes 不可用: {import_error}")
    except Exception as e:
        print(f"[PPT2PDF] WPS 阶段异常: {e}")
    finally:
        try:
            import pythoncom
            pythoncom.CoUninitialize()
        except Exception:
            pass
        _kill_wps()

    # ---------- 2) Microsoft PowerPoint ----------
    ppt_pdf = _ppt_to_pdf_powerpoint_com(file_path, pdf_path)
    if ppt_pdf and os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
        return pdf_path

    print("[PPT2PDF] WPS 与 PowerPoint COM 均失败（未安装或无法调用）")
    return None


def generate_pdf_thumbnail(file_path, output_path=None, page_num=0, scale=1.5):
    """
    生成PDF文件第一页的缩略图预览
    
    Args:
        file_path: PDF文件路径
        output_path: 输出图片路径，默认为 uploads/thumbnails/{filename}_thumb.jpg
        page_num: 要转换的页码（从0开始），默认第一页
        scale: 渲染缩放比例，默认1.5
    
    Returns:
        生成的缩略图路径，失败返回None
    """
    import os
    
    # 确保文件路径是绝对路径
    if not os.path.isabs(file_path):
        file_path = os.path.join(os.getcwd(), file_path)
    
    # 默认输出路径
    if output_path is None:
        base_dir = os.path.dirname(file_path)
        filename = os.path.basename(file_path)
        name_without_ext = os.path.splitext(filename)[0]
        thumb_dir = os.path.join(base_dir, 'thumbnails')
        output_path = os.path.join(thumb_dir, f"{name_without_ext}_thumb.jpg")
    elif not os.path.isabs(output_path):
        output_path = os.path.join(os.getcwd(), output_path)
    
    # 确保输出目录存在
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
        print(f"[Thumbnail] 创建目录: {output_dir}")
    
    # 检查文件是否存在
    if not os.path.exists(file_path):
        print(f"[Thumbnail] 文件不存在: {file_path}")
        return None
    
    # 尝试多种方法生成缩略图
    # 方法1: pypdfium2
    try:
        import pypdfium2 as pdfium
        
        doc = pdfium.PdfDocument(file_path)
        
        if page_num >= len(doc):
            print(f"[Thumbnail] 页码 {page_num} 超出范围，PDF共 {len(doc)} 页")
            return None
        
        page = doc[page_num]
        pil_img = page.render(scale=scale).to_pil()
        
        if pil_img.mode == 'RGBA':
            pil_img = pil_img.convert('RGB')
        
        pil_img.save(output_path, 'JPEG', quality=85)
        print(f"[Thumbnail] 使用pypdfium2生成缩略图成功: {output_path}")
        return output_path
        
    except ImportError:
        print("[Thumbnail] pypdfium2 未安装，尝试其他方法...")
    except Exception as e:
        print(f"[Thumbnail] pypdfium2 失败: {e}")
    
    # 方法2: PyMuPDF (fitz)
    try:
        import fitz
        
        doc = fitz.open(file_path)
        
        if page_num >= len(doc):
            print(f"[Thumbnail] 页码 {page_num} 超出范围")
            return None
        
        page = doc[page_num]
        # 渲染页面为图片
        mat = fitz.Matrix(scale, scale)
        pix = page.get_pixmap(matrix=mat)
        
        # 转换为PIL Image
        from PIL import Image
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img.save(output_path, 'JPEG', quality=85)
        
        doc.close()
        print(f"[Thumbnail] 使用PyMuPDF生成缩略图成功: {output_path}")
        return output_path
        
    except ImportError:
        print("[Thumbnail] PyMuPDF 未安装，尝试其他方法...")
    except Exception as e:
        print(f"[Thumbnail] PyMuPDF 失败: {e}")
    
    # 方法3: pdf2image
    try:
        from pdf2image import convert_from_path
        
        images = convert_from_path(file_path, first_page=page_num+1, last_page=page_num+1, dpi=150)
        
        if images:
            from PIL import Image
            img = images[0]
            width = int(img.width * scale * 0.5)
            height = int(img.height * scale * 0.5)
            img = img.resize((width, height), Image.Resampling.LANCZOS)
            img.save(output_path, 'JPEG', quality=85)
            print(f"[Thumbnail] 使用pdf2image生成缩略图成功: {output_path}")
            return output_path
            
    except ImportError:
        print("[Thumbnail] pdf2image 未安装")
    except Exception as e:
        print(f"[Thumbnail] pdf2image 失败: {e}")
    
    print("[Thumbnail] 所有方法都失败，请安装 pypdfium2 或 PyMuPDF")
    return None


def get_or_create_thumbnail(file_path, generate=True):
    """
    获取或创建课件的缩略图。

    Args:
        file_path: 课件文件路径（PDF或PPT），绝对或相对路径均可
        generate:  True（默认）= 缩略图不存在时自动生成；
                   False = 只读缓存，不存在直接返回 None（适合列表接口快速响应）

    Returns:
        缩略图相对路径（相对于 uploads 目录，用于前端访问），失败返回 None
    """
    import os

    try:
        if not os.path.isabs(file_path):
            file_path = os.path.join(os.getcwd(), file_path)

        base_dir = os.path.dirname(file_path)
        filename = os.path.basename(file_path)
        name_without_ext = os.path.splitext(filename)[0]
        ext = os.path.splitext(filename)[1].lower()

        thumb_dir = os.path.join(base_dir, 'thumbnails')
        thumb_rel = os.path.join('thumbnails', f"{name_without_ext}_thumb.jpg")
        thumb_abs = os.path.join(thumb_dir, f"{name_without_ext}_thumb.jpg")

        # 缩略图已存在 —— 直接返回，无 I/O 开销
        if os.path.exists(thumb_abs):
            return thumb_rel

        # 仅查询模式：不触发生成
        if not generate:
            return None

        print(f"[Thumbnail] 处理文件: {file_path}")

        if not os.path.exists(file_path):
            print(f"[Thumbnail] 文件不存在: {file_path}")
            return None

        # 根据文件类型生成缩略图
        if ext == '.pdf':
            result = generate_pdf_thumbnail(file_path, thumb_abs)
            if result:
                return thumb_rel

        elif ext in ['.ppt', '.pptx']:
            print(f"[Thumbnail] PPT文件，先转换为PDF...")
            pdf_path = ppt_to_pdf(file_path)
            if pdf_path and os.path.exists(pdf_path):
                result = generate_pdf_thumbnail(pdf_path, thumb_abs)
                try:
                    os.remove(pdf_path)
                except Exception:
                    pass
                if result:
                    return thumb_rel

        print(f"[Thumbnail] 无法生成缩略图: {filename}")
        return None

    except Exception as e:
        print(f"[Thumbnail] 获取缩略图失败: {e}")
        import traceback
        traceback.print_exc()
        return None